from pptx import Presentation
import pandas as pd
import docx

def read_file(current_model):
    file_folder = "./file/"
    print(current_model + "：请输入文件名", end="\n")
    filename = file_folder + input()
    if filename.endswith('.docx') or filename.endswith('.doc'):
        doc = docx.Document(filename)
        return ' '.join([paragraph.text for paragraph in doc.paragraphs])
    elif filename.endswith('.txt'):
        with open(filename, "r", encoding="utf-8") as f:
            return f.read()
    elif filename.endswith('.xlsx') or filename.endswith('.xls'):
        df = pd.read_excel(filename)
        return df.to_string(index=False)
    elif filename.endswith('.pptx'):
        prs = Presentation(filename)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return ' '.join(text)
    else:
        print("无效的文件名")
        return ""